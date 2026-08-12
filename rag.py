import os
import re
import json
import time
import shutil
import hashlib
import zipfile
import pandas as pd
from pypdf import PdfReader
from google import genai

# chromadb import แบบ lazy (เฉพาะตอนใช้จริงใน get_collection/delete_topic) เพราะเป็นแพ็กเกจหนักมาก
# (onnxruntime, tokenizers, hnswlib, duckdb ฯลฯ) — db.py (ที่ใช้กับ Streamlit Cloud) import จาก
# ไฟล์นี้แค่ read_file/chunk_text ที่ไม่เกี่ยวกับ ChromaDB เลย ไม่อยากบังคับติดตั้ง chromadb
# ทั้งก้อนโดยไม่จำเป็นตอน deploy

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
KNOWLEDGE_DIR = os.path.join(BASE_DIR, "knowledge")
CHROMA_DIR = os.path.join(BASE_DIR, "chroma_db")
os.makedirs(KNOWLEDGE_DIR, exist_ok=True)

TOPICS_META_PATH = os.path.join(KNOWLEDGE_DIR, "_topics.json")
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".webp"}
SUPPORTED_EXTS = {
    ".xlsx", ".xls", ".pdf", ".txt", ".md", ".csv", ".json", ".html", ".htm",
    ".docx", ".pptx",
} | IMAGE_EXTS
ZIP_EXT = ".zip"  # ไม่รวมใน SUPPORTED_EXTS เพราะต้องแตกไฟล์ก่อนแทนที่จะอ่านตรงๆ

# client สำหรับ embedding (ใช้ key เดียวกับ app)
_genai_client = None


def init_client(api_key):
    global _genai_client
    _genai_client = genai.Client(api_key=api_key)


EMBED_BATCH_SIZE = 100  # Gemini embedContent รับได้สูงสุด 100 รายการต่อ batch


def embed(texts):
    """แปลงข้อความเป็น vector ด้วย Gemini embedding (แบ่ง batch ถ้าเกิน 100 รายการ)"""
    embeddings = []
    for i in range(0, len(texts), EMBED_BATCH_SIZE):
        batch = texts[i:i + EMBED_BATCH_SIZE]
        result = _genai_client.models.embed_content(
            model="gemini-embedding-001",
            contents=batch,
        )
        embeddings.extend(e.values for e in result.embeddings)
    return embeddings


# ---- อ่านไฟล์แต่ละชนิด ----
def read_file(path):
    """อ่านไฟล์ตามนามสกุล คืน text"""
    ext = os.path.splitext(path)[1].lower()
    if ext in (".xlsx", ".xls"):
        df = pd.read_excel(path)
        return df.to_string(index=False)
    elif ext == ".pdf":
        reader = PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext in (".txt", ".md", ".csv", ".json"):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext in (".html", ".htm"):
        return _read_html(path)
    elif ext == ".docx":
        return _read_docx(path)
    elif ext == ".pptx":
        return _read_pptx(path)
    elif ext in IMAGE_EXTS:
        return _read_image(path)
    return ""


def _read_html(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(html, "html.parser").get_text("\n", strip=True)
    except ImportError:
        return re.sub(r"<[^>]+>", " ", html)


def _read_docx(path):
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("อ่านไฟล์ .docx ไม่ได้ ต้องติดตั้งก่อน: pip install python-docx")

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _read_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("อ่านไฟล์ .pptx ไม่ได้ ต้องติดตั้งก่อน: pip install python-pptx")

    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _read_image(path):
    """อ่านตัวหนังสือในรูปด้วย OCR ในเครื่อง (Tesseract) ไม่เรียก API ภายนอก"""
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        raise RuntimeError(
            "อ่านรูปภาพไม่ได้ ต้องติดตั้งก่อน: pip install pillow pytesseract "
            "และติดตั้งโปรแกรม Tesseract OCR ในเครื่อง "
            "(เช่น sudo apt install tesseract-ocr tesseract-ocr-tha)"
        )

    try:
        return pytesseract.image_to_string(Image.open(path), lang="tha+eng").strip()
    except pytesseract.TesseractNotFoundError:
        raise RuntimeError(
            "ไม่พบโปรแกรม Tesseract OCR ในเครื่อง ติดตั้งก่อนด้วย "
            "sudo apt install tesseract-ocr tesseract-ocr-tha"
        )


def chunk_text(text, size=800, overlap=100):
    """ตัด text เป็นชิ้นเล็กๆ (chunk) พร้อม overlap กันเนื้อหาขาด"""
    chunks = []
    start = 0
    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap
    return [c.strip() for c in chunks if c.strip()]


# ---- จัดการหัวข้อความรู้ (topic = โฟลเดอร์เฉพาะเรื่อง) ----
def _load_topics_meta():
    if not os.path.exists(TOPICS_META_PATH):
        return {}
    with open(TOPICS_META_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def _save_topics_meta(meta):
    with open(TOPICS_META_PATH, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)


def _sanitize_folder_name(name):
    """กันอักขระที่ใช้เป็นชื่อโฟลเดอร์ไม่ได้ (Windows/Unix)"""
    name = re.sub(r'[\\/:*?"<>|]', "_", name.strip())
    name = re.sub(r"\s+", " ", name)
    return name[:80]


def create_topic(display_name):
    """สร้างหัวข้อความรู้ใหม่ พร้อมโฟลเดอร์เฉพาะเรื่องนั้นใน knowledge/ ถ้าชื่อซ้ำจะคืน slug เดิม"""
    display_name = display_name.strip()
    if not display_name:
        raise ValueError("ต้องตั้งชื่อหัวข้อ")

    meta = _load_topics_meta()
    for slug, info in meta.items():
        if info["name"].strip().lower() == display_name.lower():
            return slug

    slug = "topic_" + hashlib.md5(display_name.encode("utf-8")).hexdigest()[:10]
    folder_name = _sanitize_folder_name(display_name) or slug
    os.makedirs(os.path.join(KNOWLEDGE_DIR, folder_name), exist_ok=True)

    meta[slug] = {
        "name": display_name,
        "folder": folder_name,
        "created_at": time.time(),
    }
    _save_topics_meta(meta)
    return slug


def list_topics():
    """คืนรายชื่อหัวข้อทั้งหมด เรียงตามชื่อ"""
    meta = _load_topics_meta()
    return [
        {"slug": slug, "name": info["name"], "folder": info["folder"]}
        for slug, info in sorted(meta.items(), key=lambda kv: kv[1]["name"])
    ]


def get_topic_dir(topic_slug):
    """คืน path โฟลเดอร์ของหัวข้อนั้น หรือ None ถ้าไม่มีหัวข้อนี้"""
    info = _load_topics_meta().get(topic_slug)
    if not info:
        return None
    return os.path.join(KNOWLEDGE_DIR, info["folder"])


def rename_topic(topic_slug, new_name):
    """แก้ชื่อที่แสดงของหัวข้อ (โฟลเดอร์/collection เดิมไม่เปลี่ยน)"""
    new_name = new_name.strip()
    if not new_name:
        raise ValueError("ต้องตั้งชื่อหัวข้อ")

    meta = _load_topics_meta()
    if topic_slug not in meta:
        raise ValueError("ไม่พบหัวข้อนี้")

    meta[topic_slug]["name"] = new_name
    _save_topics_meta(meta)


def delete_topic(topic_slug):
    """ลบหัวข้อทั้งหมด: โฟลเดอร์ไฟล์ + collection ใน ChromaDB + metadata"""
    info = _load_topics_meta().get(topic_slug)
    if not info:
        return

    topic_dir = os.path.join(KNOWLEDGE_DIR, info["folder"])
    if os.path.isdir(topic_dir):
        shutil.rmtree(topic_dir)

    try:
        import chromadb
        chromadb.PersistentClient(path=CHROMA_DIR).delete_collection(name=topic_slug)
    except Exception:
        pass  # ไม่มี collection อยู่แล้วก็ข้ามไป

    meta = _load_topics_meta()
    meta.pop(topic_slug, None)
    _save_topics_meta(meta)


# ---- จัดการ collection แต่ละหัวข้อ/โหมด ----
def get_collection(mode_key):
    """แต่ละหัวข้อ/โหมดมี collection แยกกัน"""
    import chromadb
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    return client.get_or_create_collection(name=mode_key)


def add_document(mode_key, path, filename):
    """อ่านไฟล์ → chunk → embed → เก็บลง ChromaDB (.zip จะถูกแตกไฟล์แล้วเพิ่มทีละไฟล์ข้างในแทน)"""
    if os.path.splitext(filename)[1].lower() == ZIP_EXT:
        return _add_zip(mode_key, path)

    text = read_file(path)
    if not text:
        return 0

    chunks = chunk_text(text)
    embeddings = embed(chunks)

    collection = get_collection(mode_key)
    ids = [f"{filename}_{i}" for i in range(len(chunks))]
    collection.add(
        ids=ids,
        documents=chunks,
        embeddings=embeddings,
        metadatas=[{"source": filename} for _ in chunks],
    )
    return len(chunks)


def delete_source_file(topic_slug, filename):
    """ลบไฟล์เดียวออกจากหัวข้อ: ลบไฟล์จริงในโฟลเดอร์ + ลบ chunk ที่เกี่ยวข้องใน ChromaDB"""
    topic_dir = get_topic_dir(topic_slug)
    if topic_dir:
        path = os.path.join(topic_dir, filename)
        if os.path.exists(path):
            os.remove(path)

    collection = get_collection(topic_slug)
    if collection.count() > 0:
        all_data = collection.get()
        ids_to_delete = [
            id_ for id_, meta in zip(all_data["ids"], all_data["metadatas"])
            if meta.get("source") == filename
        ]
        if ids_to_delete:
            collection.delete(ids=ids_to_delete)


def replace_source_file(topic_slug, path, filename):
    """แทนที่ไฟล์เดิม (ที่ path ถูกเขียนทับ/อัปโหลดใหม่แล้ว) ด้วยการลบ chunk เก่าแล้ว re-embed ใหม่ทั้งไฟล์"""
    delete_source_file(topic_slug, filename)
    return add_document(topic_slug, path, filename)


def _add_zip(topic_slug, zip_path):
    """แตกไฟล์ในซิปลงโฟลเดอร์ของหัวข้อเดียวกัน แล้วเพิ่มทุกไฟล์ที่รองรับข้างใน (กันเพิ่มซ้ำเหมือนไฟล์ทั่วไป)"""
    topic_dir = os.path.dirname(zip_path)
    collection = get_collection(topic_slug)
    total = 0
    with zipfile.ZipFile(zip_path, "r") as zf:
        for member in zf.namelist():
            if member.endswith("/"):
                continue
            # ใช้แค่ basename กันไฟล์ในซิปเขียนออกนอกโฟลเดอร์หัวข้อ (zip-slip)
            inner_name = os.path.basename(member)
            inner_ext = os.path.splitext(inner_name)[1].lower()
            if not inner_name or inner_ext not in SUPPORTED_EXTS:
                continue

            extract_path = os.path.join(topic_dir, inner_name)
            with zf.open(member) as src, open(extract_path, "wb") as dst:
                dst.write(src.read())

            existing = collection.get(ids=[f"{inner_name}_0"])
            if not existing["ids"]:
                total += add_document(topic_slug, extract_path, inner_name)
    return total


def load_folder(topic_slug):
    """โหลดไฟล์ทั้งหมดในโฟลเดอร์เฉพาะของหัวข้อนี้ที่ยังไม่เคยเพิ่ม (รองรับ .zip ด้วย)"""
    topic_dir = get_topic_dir(topic_slug)
    if not topic_dir or not os.path.isdir(topic_dir):
        return 0

    collection = get_collection(topic_slug)
    count = 0
    for fname in os.listdir(topic_dir):
        ext = os.path.splitext(fname)[1].lower()
        path = os.path.join(topic_dir, fname)
        if ext == ZIP_EXT:
            count += _add_zip(topic_slug, path)
        elif ext in SUPPORTED_EXTS:
            # กันเพิ่มซ้ำ: เช็คว่ามี id นี้แล้วไหม
            existing = collection.get(ids=[f"{fname}_0"])
            if not existing["ids"]:
                count += add_document(topic_slug, path, fname)
    return count


def search(mode_key, query, top_k=4):
    """ค้นหา chunk ที่เกี่ยวกับ query มากที่สุด"""
    collection = get_collection(mode_key)
    if collection.count() == 0:
        return ""
    query_emb = embed([query])[0]
    results = collection.query(query_embeddings=[query_emb], n_results=top_k)
    docs = results["documents"][0] if results["documents"] else []
    return "\n\n---\n\n".join(docs)


def list_sources(mode_key):
    """คืนรายชื่อไฟล์ที่อยู่ใน collection"""
    collection = get_collection(mode_key)
    if collection.count() == 0:
        return []
    all_meta = collection.get()["metadatas"]
    return sorted(set(m["source"] for m in all_meta))